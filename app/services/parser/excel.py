import pandas as pd
from io import BytesIO
from PIL import Image
import pytesseract
from pptx import Presentation
import zipfile
import os
import csv
import io
import time

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        # Read the text content from the bytes
        with BytesIO(file_bytes) as txt_file:
            content = txt_file.read().decode('utf-8', errors='ignore')
        return content
    except Exception as e:
        return f"Error reading TXT file: {str(e)}"



def extract_text_from_excel_bytes(data: bytes) -> str:
    excel_file = BytesIO(data)
    df = pd.read_excel(excel_file)
    return df.to_string(index=False)


def extract_text_from_image_bytes(data: bytes) -> str:
    image = Image.open(BytesIO(data))
    return pytesseract.image_to_string(image)



def extract_text_from_pptx_with_ocr(data: bytes) -> str:
    text = []

    # Load presentation
    presentation = Presentation(BytesIO(data))

    # Extract visible text from shapes
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    # Extract images from pptx as a zip file
    zip_data = zipfile.ZipFile(BytesIO(data))
    image_files = [f for f in zip_data.namelist() if f.startswith("ppt/media/")]

    for image_name in image_files:
        with zip_data.open(image_name) as image_file:
            try:
                image = Image.open(image_file)
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    text.append(ocr_text)
            except Exception as e:
                print(f"OCR error on {image_name}: {e}")

    return "\n".join(text)

def extract_text_from_csv_bytes(file_bytes: bytes) -> str:
    text = ""
    with io.StringIO(file_bytes.decode('utf-8', errors='ignore')) as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            text += ", ".join(row) + "\n"
    return text.strip()



def extract_text_from_nested_zip(zip_bytes, max_depth=10):
    start_time = time.time()
    result = []

    def process_zip(zip_bytes, depth):
        nonlocal start_time
        if depth > max_depth:
            result.append("Maximum zip nesting level reached.")
            return

        if time.time() - start_time > 7:
            result.append("Timeout: Nested zip extraction.")
            return

        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            for name in zf.namelist():
                with zf.open(name) as file:
                    if name.endswith(".txt"):
                        result.append(f"\n[Text file: {name}]\n" + file.read().decode(errors="ignore"))
                    elif name.endswith(".csv"):
                        result.append(f"\n[CSV file: {name}]\n" + file.read().decode(errors="ignore"))
                    elif name.endswith(".zip"):
                        try:
                            extract_text = file.read()
                            process_zip(extract_text, depth + 1)
                        except Exception as e:
                            result.append(f"[ERROR] Could not process nested zip '{name}': {str(e)}")

    try:
        process_zip(zip_bytes, 0)
    except Exception as e:
        result.append(f"[ERROR] Failed to process zip: {str(e)}")

    return "\n".join(result) if result else "[INFO] No readable content found."
